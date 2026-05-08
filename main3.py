import os
import json
import re
import uuid
import warnings
from datetime import datetime
from typing import List, Optional, Any, Dict, Tuple

os.environ.setdefault("KMP_DUPLICATE_LIB_OK", "TRUE")
warnings.filterwarnings("ignore", category=DeprecationWarning)
warnings.filterwarnings("ignore", message=".*NotOpenSSLWarning.*")

import base64
import fitz
from pydantic import Field
from langchain.llms.base import LLM
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from groq import Groq
from dotenv import load_dotenv
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS

load_dotenv()

_api_key = os.getenv("GROQ_API_KEY")
if not _api_key:
    print("\n" + "="*60)
    print("  ERROR: GROQ_API_KEY is not set!")
    print("="*60)
    print("  Fix: create a file called  .env  in this folder with:")
    print()
    print("    GROQ_API_KEY=gsk_your-key-here")
    print()
    print("  Or export it in your terminal before running:")
    print("    export GROQ_API_KEY=gsk_your-key-here")
    print("="*60 + "\n")

def get_api_key():
    key = os.getenv("GROQ_API_KEY")
    if not key:
        raise ValueError(
            "GROQ_API_KEY is not set. Add it to your .env file: "
            "GROQ_API_KEY=gsk_your-key-here"
        )
    return key

app = Flask(__name__, static_folder=".")
CORS(app, expose_headers=["Content-Disposition"])

RFP_PROJECTS_ROOT = "rfp_projects"
PROJECTS_FILE     = "projects.json"
os.makedirs(RFP_PROJECTS_ROOT, exist_ok=True)


def _next_project_number() -> int:
    existing = [
        d for d in os.listdir(RFP_PROJECTS_ROOT)
        if os.path.isdir(os.path.join(RFP_PROJECTS_ROOT, d))
    ]
    nums = []
    for name in existing:
        parts = name.split("_", 1)
        try:
            nums.append(int(parts[0]))
        except ValueError:
            pass
    return max(nums, default=0) + 1


def _make_project_folder(doc_title: str) -> Tuple[str, str]:
    num        = _next_project_number()
    safe_title = re.sub(r"[^a-zA-Z0-9]+", "_", (doc_title or "Untitled").strip())[:50].strip("_")
    folder_name = f"{num:03d}_{safe_title}"
    folder_path = os.path.join(RFP_PROJECTS_ROOT, folder_name)
    os.makedirs(folder_path, exist_ok=True)
    os.makedirs(os.path.join(folder_path, "vectorstore"), exist_ok=True)
    return folder_path, folder_name


def _project_folder(project: dict) -> str:
    folder = project.get("folder", "")
    if folder and os.path.exists(folder):
        return folder
    return RFP_PROJECTS_ROOT


# ─── LLM (Groq) ───────────────────────────────────────────────────────────────
class GroqLLM(LLM):
    client: Any = Field(default=None)
    model_name: str = "llama-3.3-70b-versatile"
    system_prompt: str = (
        "You are an expert RFP document analyst. "
        "Extract structured information from Request for Proposal documents. "
        "Always respond with valid JSON only — no markdown, no explanation."
    )

    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        if self.client is None:
            object.__setattr__(self, 'client', Groq(api_key=get_api_key()))

    def _call(self, prompt: str, stop: Optional[List[str]] = None) -> str:
        response = self.client.chat.completions.create(
            model=self.model_name,
            messages=[
                {"role": "system", "content": self.system_prompt},
                {"role": "user",   "content": prompt},
            ],
            temperature=0.2,
            max_tokens=4000,
            stop=stop,
        )
        return response.choices[0].message.content

    @property
    def _llm_type(self) -> str:
        return "groq"


def get_embeddings():
    """Return HuggingFace sentence-transformer embeddings (replaces OpenAIEmbeddings)."""
    return HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2",
        model_kwargs={"device": "cpu"},
        encode_kwargs={"normalize_embeddings": True},
    )


# ─── PROJECT STORE ─────────────────────────────────────────────────────────────
def load_projects() -> Dict:
    if os.path.exists(PROJECTS_FILE):
        with open(PROJECTS_FILE, "r") as f:
            return json.load(f)
    return {}


def save_projects(projects: Dict):
    with open(PROJECTS_FILE, "w") as f:
        json.dump(projects, f, indent=2)


def save_intelligence_report(project_id: str, doc_title: str, intelligence: dict,
                             folder: str = "") -> str:
    dest = folder if folder else RFP_PROJECTS_ROOT
    report_path = os.path.join(dest, "intelligence.json")
    payload = {
        "project_id":   project_id,
        "doc_title":    doc_title,
        "extracted_at": datetime.now().isoformat(),
        **intelligence,
    }
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(payload, f, indent=2, ensure_ascii=False)
    print(f"[INFO] Intelligence report saved → {report_path}")
    return report_path


# ─── DOCUMENT PROCESSING ──────────────────────────────────────────────────────
_MIN_IMAGE_PX  = 80
_MAX_IMAGES_PG = 4
_MAX_TABLES_PG = 6


def _extract_page_tables(page) -> list:
    try:
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
    except Exception:
        return []

    words = []
    for block in blocks:
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                t = span.get("text", "").strip()
                if not t:
                    continue
                x0, y0, x1, y1 = span["bbox"]
                words.append({"x": (x0+x1)/2, "y": (y0+y1)/2, "x0": x0, "text": t})

    if len(words) < 6:
        return []

    Y_TOL = 6
    rows_dict: dict = {}
    for w in words:
        placed = False
        for yk in list(rows_dict.keys()):
            if abs(w["y"] - yk) <= Y_TOL:
                rows_dict[yk].append(w)
                placed = True
                break
        if not placed:
            rows_dict[w["y"]] = [w]

    sorted_rows = []
    for yk in sorted(rows_dict.keys()):
        row_words = sorted(rows_dict[yk], key=lambda w: w["x0"])
        sorted_rows.append(row_words)

    tables = []
    in_table = False
    table_rows: list = []
    for row in sorted_rows:
        if len(row) >= 2:
            table_rows.append(row)
            in_table = True
        else:
            if in_table and len(table_rows) >= 3:
                tables.append(table_rows)
            table_rows = []
            in_table = False
    if in_table and len(table_rows) >= 3:
        tables.append(table_rows)

    if not tables:
        return []

    md_tables = []
    for tbl in tables[:_MAX_TABLES_PG]:
        header_row = tbl[0]
        col_xs = [w["x"] for w in header_row]

        def assign_col(word_x):
            return min(range(len(col_xs)), key=lambda i: abs(col_xs[i] - word_x))

        md_rows = []
        for row_words in tbl:
            cells = [""] * len(col_xs)
            for w in row_words:
                col = assign_col(w["x"])
                cells[col] = (cells[col] + " " + w["text"]).strip()
            md_rows.append("| " + " | ".join(cells) + " |")

        if len(md_rows) >= 2:
            sep = "| " + " | ".join(["---"] * len(col_xs)) + " |"
            md_table = "\n".join([md_rows[0], sep] + md_rows[1:])
            md_tables.append(md_table)

    return md_tables


def _describe_page_images(page, page_num: int) -> list:
    """
    Image description via Groq vision (llama-3.2-90b-vision-preview).
    Falls back gracefully if vision model is unavailable.
    """
    descriptions = []
    img_list = page.get_images(full=True)
    if not img_list:
        return descriptions

    doc = page.parent
    groq_client = Groq(api_key=get_api_key())
    seen_xrefs = set()
    count = 0

    for img_info in img_list:
        if count >= _MAX_IMAGES_PG:
            break
        xref = img_info[0]
        if xref in seen_xrefs:
            continue
        seen_xrefs.add(xref)

        try:
            pix = fitz.Pixmap(doc, xref)
            if pix.width < _MIN_IMAGE_PX or pix.height < _MIN_IMAGE_PX:
                pix = None
                continue
            if pix.n > 4:
                pix = fitz.Pixmap(fitz.csRGB, pix)

            img_bytes = pix.tobytes("png")
            img_b64   = base64.b64encode(img_bytes).decode("utf-8")
            pix       = None

            response = groq_client.chat.completions.create(
                model="llama-3.2-90b-vision-preview",
                max_tokens=300,
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "You are analysing an image extracted from an RFP / business document. "
                                "Describe what this image shows in 1-3 sentences, focusing on: "
                                "any numbers, percentages, dates, names, titles, or key findings visible. "
                                "If it is a chart or graph, state the type, axes labels, and the main trend or value. "
                                "If it is a table rendered as an image, transcribe the key data. "
                                "If it is a logo, diagram, or decorative element with no data, say so briefly."
                            ),
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/png;base64,{img_b64}"},
                        },
                    ],
                }],
            )
            desc = response.choices[0].message.content.strip()
            if desc:
                descriptions.append(desc)
                count += 1

        except Exception as img_err:
            print(f"[WARN] Page {page_num+1} image description failed: {img_err}")
            continue

    return descriptions


def extract_text_with_pages(file_path: str, describe_images: bool = True) -> List[Dict]:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="replace") as fh:
            raw = fh.read()
        if "<physical_index_1>" in raw:
            return parse_structured_txt(raw)
        return [{"page": 1, "text": raw, "images": [], "tables": []}]

    if ext != ".pdf":
        raise ValueError(f"Unsupported file type '{ext}'. Only .pdf and .txt accepted.")

    try:
        doc = fitz.open(file_path)
    except Exception as e:
        raise RuntimeError(f"PyMuPDF could not open '{file_path}': {e}")

    pages = []
    total = len(doc)

    for page_num in range(total):
        page = doc[page_num]
        pn   = page_num + 1
        text = page.get_text("text", flags=fitz.TEXT_PRESERVE_WHITESPACE)
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        tables = _extract_page_tables(page)
        images = []
        if describe_images:
            try:
                images = _describe_page_images(page, page_num)
            except Exception as ie:
                print(f"[WARN] Page {pn} image extraction failed: {ie}")

        if text or images or tables:
            pages.append({"page": pn, "text": text, "images": images, "tables": tables})

        print(f"[INFO] Page {pn}/{total}: {len(text)} chars text, {len(tables)} table(s), {len(images)} image(s)")

    doc.close()
    return pages


def build_structured_txt(pages: List[Dict]) -> str:
    parts = []
    for p in pages:
        n      = p["page"]
        chunks = ["[TEXT]", p["text"].strip()] if p.get("text", "").strip() else []
        for ti, tbl in enumerate(p.get("tables", []), 1):
            chunks.append(f"\n[TABLE {ti}]:\n{tbl}")
        for ii, desc in enumerate(p.get("images", []), 1):
            chunks.append(f"\n[IMAGE {ii}]: {desc}")
        body = "\n\n".join(chunks).strip()
        if body:
            parts.append(f"<physical_index_{n}>\n{body}\n</physical_index_{n}>")
    return "\n\n".join(parts)


def save_structured_txt(pages: List[Dict], folder: str) -> str:
    os.makedirs(folder, exist_ok=True)
    txt_path = os.path.join(folder, "extracted.txt")
    with open(txt_path, "w", encoding="utf-8") as fh:
        fh.write(build_structured_txt(pages))
    n_img = sum(len(p.get("images", [])) for p in pages)
    n_tbl = sum(len(p.get("tables", [])) for p in pages)
    print(f"[INFO] extracted.txt saved → {txt_path}  ({len(pages)} pages, {n_tbl} tables, {n_img} images described)")
    return txt_path


def parse_structured_txt(raw: str) -> List[Dict]:
    page_pattern = re.compile(r"<physical_index_(\d+)>\s*(.*?)\s*</physical_index_\1>", re.DOTALL)
    tbl_pattern  = re.compile(r"\[TABLE \d+\]:\s*((?:\|[^\n]*\n?)+)", re.DOTALL)
    img_pattern  = re.compile(r"\[IMAGE \d+\]:\s*(.+?)(?=\n\[|\Z)", re.DOTALL)

    pages = []
    for m in page_pattern.finditer(raw):
        body = m.group(2).strip()
        text_match = re.search(r"\[TEXT\](.*?)(?=\n\[TABLE|\n\[IMAGE|\Z)", body, re.DOTALL)
        text = text_match.group(1).strip() if text_match else ""
        if not text and not re.search(r"\[TABLE|\[IMAGE", body):
            text = body
        tables = [t.strip() for t in tbl_pattern.findall(body)]
        images = [i.strip() for i in img_pattern.findall(body)]
        if text or tables or images:
            pages.append({"page": int(m.group(1)), "text": text, "tables": tables, "images": images})

    if not pages:
        pages = [{"page": 1, "text": raw.strip(), "images": [], "tables": []}]
    return pages


def load_structured_txt(folder: str) -> List[Dict]:
    txt_path = os.path.join(folder, "extracted.txt")
    if not os.path.exists(txt_path):
        raise FileNotFoundError(f"extracted.txt not found in '{folder}'.")
    with open(txt_path, "r", encoding="utf-8") as fh:
        return parse_structured_txt(fh.read())


def get_full_text_from_pages(pages: List[Dict]) -> str:
    parts = []
    for p in pages:
        section = [f"[Page {p['page']}]"]
        if p.get("text", "").strip():
            section.append(p["text"].strip())
        for ti, tbl in enumerate(p.get("tables", []), 1):
            section.append(f"[TABLE {ti}]\n{tbl}")
        for ii, desc in enumerate(p.get("images", []), 1):
            section.append(f"[IMAGE {ii}]: {desc}")
        parts.append("\n\n".join(section))
    return "\n\n".join(parts)


def get_pages_for_project(project: dict) -> List[Dict]:
    folder   = _project_folder(project)
    txt_path = os.path.join(folder, "extracted.txt")
    if os.path.exists(txt_path):
        try:
            return load_structured_txt(folder)
        except Exception as e:
            print(f"[WARN] Could not load extracted.txt: {e}. Re-extracting.")
    file_path = project.get("file_path", "")
    if file_path and os.path.exists(file_path):
        pages = extract_text_with_pages(file_path)
        save_structured_txt(pages, folder)
        return pages
    raise FileNotFoundError(f"No extracted.txt or source file for project '{project.get('id', '?')}'.")


def extract_text(file_path: str) -> str:
    return get_full_text_from_pages(extract_text_with_pages(file_path))


def build_vectorstore(project_id: str, pages: List[Dict], doc_meta: Dict):
    splitter   = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=120)
    embeddings = get_embeddings()
    texts, metadatas = [], []
    for page_info in pages:
        chunks = splitter.split_text(page_info["text"])
        for chunk in chunks:
            texts.append(chunk)
            metadatas.append({
                "document_id":      project_id,
                "document_name":    doc_meta.get("doc_title", ""),
                "company_name":     doc_meta.get("company_name", ""),
                "page_number":      page_info["page"],
                "project_category": doc_meta.get("category", "General"),
                "project_status":   doc_meta.get("status", "active"),
            })
    vs = FAISS.from_texts(texts, embeddings, metadatas=metadatas)
    projects  = load_projects()
    proj      = projects.get(project_id, {})
    folder    = _project_folder(proj)
    store_path = os.path.join(folder, "vectorstore")
    os.makedirs(store_path, exist_ok=True)
    vs.save_local(store_path)
    return store_path


def load_vectorstore(project_id: str):
    embeddings = get_embeddings()
    projects   = load_projects()
    proj       = projects.get(project_id, {})
    folder     = _project_folder(proj)
    store_path = os.path.join(folder, "vectorstore")
    if not os.path.exists(store_path):
        return None
    return FAISS.load_local(store_path, embeddings, allow_dangerous_deserialization=True)


# ─── LLM HELPERS ──────────────────────────────────────────────────────────────
def parse_json_response(raw: str) -> dict:
    cleaned = re.sub(r"```(?:json)?", "", raw).strip().rstrip("`").strip()
    match   = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if match:
        cleaned = match.group(0)
    return json.loads(cleaned)


def groq_chat(prompt: str, model: str = "llama-3.3-70b-versatile",
              system: str = "You are an expert RFP document analyst. Always respond with valid JSON only — no markdown, no explanation.",
              max_tokens: int = 4000, temperature: float = 0.2) -> str:
    """Thin wrapper around Groq chat completions."""
    client = Groq(api_key=get_api_key())
    resp   = client.chat.completions.create(
        model=model,
        messages=[{"role": "system", "content": system}, {"role": "user", "content": prompt}],
        temperature=temperature,
        max_tokens=max_tokens,
    )
    return resp.choices[0].message.content


def extract_doc_meta(llm: GroqLLM, document_text: str) -> dict:
    prompt = f"""
Analyze the following RFP document and extract key metadata.

Return ONLY valid JSON in this exact format:
{{
  "doc_title": "Full official name of the RFP document",
  "company_name": "Name of the issuing organization or company",
  "rfp_number": "RFP/tender reference number if present, else null",
  "category": "One of: Medical, Technical, Retail, Construction, IT, Finance, Legal, Education, Government, Other",
  "status": "active"
}}

Document (first 3000 chars):
{document_text[:3000]}
"""
    raw = llm.invoke(prompt)
    return parse_json_response(raw)


def extract_sections(llm: GroqLLM, document_text: str) -> list:
    prompt = f"""
Analyze the following RFP document and extract all major sections.

Rules:
- Identify clear structural headings only (not sub-bullets or minor notes).
- Give each section a short slug id (snake_case, no spaces).
- Give each section a short one-line description (max 8 words).
- Return ONLY valid JSON in this exact format:

{{
  "sections": [
    {{"id": "scope_of_work", "title": "Scope of Work", "desc": "Vendor deliverables and responsibilities"}},
    {{"id": "budget", "title": "Budget", "desc": "Financial allocation and constraints"}}
  ]
}}

Document (first 14000 chars):
{document_text[:14000]}
"""
    raw  = llm.invoke(prompt)
    data = parse_json_response(raw)
    return data.get("sections", [])


def generate_detailed_briefs(llm: GroqLLM, document_text: str, sections: list) -> dict:
    section_titles = "\n".join(f"- {s['title']} (id: {s['id']})" for s in sections)
    prompt = f"""
You are given an RFP document and a list of sections identified in it.
For EACH section, extract a DETAILED BRIEF — thorough, faithful extraction of what that section actually says.
Write it as 4-8 sentences of dense, information-rich prose.

Sections to extract:
{section_titles}

Return ONLY valid JSON:
{{
  "briefs": [
    {{"id": "scope_of_work", "brief": "Full detailed content of this section..."}}
  ]
}}

Full Document:
{document_text[:18000]}
"""
    raw  = llm.invoke(prompt)
    data = parse_json_response(raw)
    return {item["id"]: item["brief"] for item in data.get("briefs", [])}


def generate_summaries(llm: GroqLLM, document_text: str, sections: list) -> dict:
    section_titles = "\n".join(f"- {s['title']} (id: {s['id']})" for s in sections)
    prompt = f"""
For EACH section below, write a concise 2-3 sentence summary based on the document content.

Sections:
{section_titles}

Return ONLY valid JSON:
{{
  "summaries": [
    {{"id": "scope_of_work", "summary": "This section covers..."}},
    {{"id": "budget", "summary": "The budget section outlines..."}}
  ]
}}

Document (first 14000 chars):
{document_text[:14000]}
"""
    raw  = llm.invoke(prompt)
    data = parse_json_response(raw)
    return {item["id"]: item["summary"] for item in data.get("summaries", [])}


def generate_keywords(llm: GroqLLM, document_text: str, sections: list, briefs: dict, summaries: dict) -> dict:
    section_info = "\n".join(f"- {s['title']} (id: {s['id']})" for s in sections)
    prompt = f"""
For each section listed below, identify 5-12 important keywords or short phrases (2-4 words max each).

Sections:
{section_info}

Return ONLY valid JSON:
{{
  "keywords": [
    {{"id": "scope_of_work", "words": ["acceptance criteria", "phased delivery", "IP ownership"]}}
  ]
}}

Document context (first 10000 chars):
{document_text[:10000]}
"""
    raw  = llm.invoke(prompt)
    data = parse_json_response(raw)
    return {item["id"]: item["words"] for item in data.get("keywords", [])}


def generate_taglines(llm: GroqLLM, document_text: str, sections: list, doc_title: str, company_name: str) -> dict:
    section_titles = "\n".join(f"- {s['title']} (id: {s['id']})" for s in sections)
    prompt = f"""
Document: "{doc_title}" by "{company_name}"

For each section below, generate 2-4 short tag labels (1-3 words each).

Sections:
{section_titles}

Return ONLY valid JSON:
{{
  "taglines": [
    {{"id": "introduction", "tags": ["Tag One", "Tag Two", "Tag Three"]}}
  ]
}}

Document (first 5000 chars):
{document_text[:5000]}
"""
    raw  = llm.invoke(prompt)
    data = parse_json_response(raw)
    return {item["id"]: item["tags"] for item in data.get("taglines", [])}


def generate_suggestions(llm: GroqLLM, document_text: str, sections: list, summaries: dict) -> dict:
    section_titles = "\n".join(
        f"- {s['title']} (id: {s['id']}): {summaries.get(s['id'], '')}" for s in sections
    )
    prompt = f"""
For each section below, provide exactly 3 short, practical, actionable suggestions.

Sections and summaries:
{section_titles}

Return ONLY valid JSON:
{{
  "suggestions": [
    {{
      "id": "scope_of_work",
      "items": [
        "Clearly define deliverable acceptance criteria.",
        "Ask for phased delivery timeline.",
        "Ensure IP ownership clauses are stated."
      ]
    }}
  ]
}}

Document context (first 8000 chars):
{document_text[:8000]}
"""
    raw  = llm.invoke(prompt)
    data = parse_json_response(raw)
    return {item["id"]: item["items"] for item in data.get("suggestions", [])}


def infer_tag(title: str) -> str:
    t = title.lower()
    if any(k in t for k in ["goal","objective","purpose"]): return "Strategy"
    if any(k in t for k in ["scope","work","deliverable"]): return "Requirements"
    if any(k in t for k in ["budget","cost","financial","price"]): return "Financial"
    if any(k in t for k in ["timeline","schedule","milestone"]): return "Schedule"
    if any(k in t for k in ["compliance","legal","regulatory"]): return "Legal"
    if any(k in t for k in ["evaluat","criteria","scoring"]): return "Evaluation"
    if any(k in t for k in ["company","about","background","intro"]): return "Context"
    if any(k in t for k in ["current","situation","existing"]): return "Analysis"
    if any(k in t for k in ["technical","architecture","system"]): return "Technical"
    return "Section"


def build_tree_data(document_text, sections, briefs, summaries, suggestions, keywords, taglines, doc_title, company_name):
    nodes = []
    for s in sections:
        nodes.append({
            "id":          s["id"],
            "title":       s["title"],
            "desc":        s.get("desc", ""),
            "tag":         infer_tag(s["title"]),
            "brief":       briefs.get(s["id"], "No detailed brief available."),
            "summary":     summaries.get(s["id"], "No summary available."),
            "suggestions": suggestions.get(s["id"], []),
            "keywords":    keywords.get(s["id"], []),
            "taglines":    taglines.get(s["id"], [])
        })
    return {"doc_title": doc_title, "company_name": company_name, "nodes": nodes}


def extract_rfp_intelligence(llm: GroqLLM, document_text: str) -> dict:
    prompt = f"""
You are an expert RFP analyst. Analyze the following RFP document and extract structured intelligence.

Return ONLY valid JSON:
{{
  "domain": "The high-level domain this RFP belongs to",
  "sector": "Specific industry sector",
  "summary": "A concise 3-5 sentence executive summary",
  "keywords": ["keyword1", "keyword2", "keyword3", "keyword4", "keyword5", "keyword6"],
  "pricing": {{
    "budget_mentioned": true,
    "estimated_value": "Exact value or inferred range",
    "pricing_model": "Fixed Price / Time & Material / Milestone-based / Not specified",
    "currency": "USD / INR / GBP / EUR / inferred",
    "notes": "Any relevant pricing notes"
  }},
  "stakeholders": [
    {{"name": "Organisation or person name", "role": "Client / Vendor / Evaluator", "contact": null}}
  ],
  "deadlines": [
    {{"event": "Proposal Submission", "date": "DD MMM YYYY or Not specified", "is_inferred": false}}
  ],
  "tech_categories": [
    {{"category": "Technology category", "specifics": ["specific tech"], "mandatory": true}}
  ]
}}

Document (first 16000 chars):
{document_text[:16000]}
"""
    raw = llm.invoke(prompt)
    return parse_json_response(raw)


# ─── FLASK ROUTES ──────────────────────────────────────────────────────────────
@app.route("/")
def serve_index():
    return send_from_directory(".", "index.html")

@app.route("/<path:filename>")
def serve_static(filename):
    return send_from_directory(".", filename)

@app.route("/api/projects", methods=["GET"])
def get_projects():
    projects = load_projects()
    return jsonify(list(projects.values()))


@app.route("/api/upload", methods=["POST"])
def upload_project():
    if not os.getenv("GROQ_API_KEY"):
        return jsonify({"error": "GROQ_API_KEY is not set."}), 500
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file           = request.files["file"]
    category       = request.form.get("category", "General")
    status         = request.form.get("status", "active")
    employee_count = int(request.form.get("employee_count", 0))
    project_id     = str(uuid.uuid4())[:8]

    import tempfile, shutil
    ext      = os.path.splitext(file.filename)[1].lower()
    tmp_path = os.path.join(tempfile.gettempdir(), f"rfp_tmp_{project_id}{ext}")
    file.save(tmp_path)

    try:
        pages     = extract_text_with_pages(tmp_path, describe_images=True)
        full_text = get_full_text_from_pages(pages)

        llm  = GroqLLM()
        meta = extract_doc_meta(llm, full_text)
        meta["category"] = category
        meta["status"]   = status

        doc_title    = meta.get("doc_title") or os.path.splitext(file.filename)[0]
        company_name = meta.get("company_name", "")

        folder_path, folder_name = _make_project_folder(doc_title)
        file_path = os.path.join(folder_path, f"original{ext}")
        shutil.move(tmp_path, file_path)
        txt_path = save_structured_txt(pages, folder_path)

        sections    = extract_sections(llm, full_text)
        briefs      = generate_detailed_briefs(llm, full_text, sections)
        summaries   = generate_summaries(llm, full_text, sections)
        keywords    = generate_keywords(llm, full_text, sections, briefs, summaries)
        taglines    = generate_taglines(llm, full_text, sections, doc_title, company_name)
        suggestions = generate_suggestions(llm, full_text, sections, summaries)

        tree_data = build_tree_data(
            full_text, sections, briefs, summaries,
            suggestions, keywords, taglines, doc_title, company_name
        )

        data_file = os.path.join(folder_path, "data.json")
        with open(data_file, "w", encoding="utf-8") as f:
            json.dump(tree_data, f, indent=2, ensure_ascii=False)

        build_vectorstore(project_id, pages, meta)

        try:
            intelligence = extract_rfp_intelligence(llm, full_text)
            save_intelligence_report(project_id, doc_title, intelligence, folder=folder_path)
        except Exception as intel_err:
            print(f"[WARN] Intelligence extraction failed: {intel_err}")

        projects = load_projects()
        projects[project_id] = {
            "id":            project_id,
            "folder":        folder_path,
            "folder_name":   folder_name,
            "doc_title":     doc_title,
            "company_name":  company_name,
            "category":      category,
            "status":        status,
            "employee_count": employee_count,
            "rfp_number":    meta.get("rfp_number"),
            "data_file":     data_file,
            "file_path":     file_path,
            "txt_path":      txt_path,
            "created_at":    datetime.now().isoformat(),
            "page_count":    len(pages),
            "section_count": len(sections),
        }
        save_projects(projects)

        return jsonify({
            "success":    True,
            "project_id": project_id,
            "folder":     folder_name,
            "data_file":  data_file,
            "project":    projects[project_id],
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/projects/<project_id>", methods=["PATCH"])
def update_project(project_id):
    projects = load_projects()
    if project_id not in projects:
        return jsonify({"error": "Project not found"}), 404
    data    = request.get_json()
    allowed = {"status", "category", "employee_count", "doc_title", "company_name"}
    for key in allowed:
        if key in data:
            projects[project_id][key] = data[key]
    save_projects(projects)
    return jsonify(projects[project_id])


@app.route("/api/projects/<project_id>", methods=["DELETE"])
def delete_project(project_id):
    projects = load_projects()
    if project_id not in projects:
        return jsonify({"error": "Project not found"}), 404
    del projects[project_id]
    save_projects(projects)
    return jsonify({"success": True})


@app.route("/api/chat", methods=["POST"])
def chat():
    try:
        body       = request.get_json()
        if not body:
            return jsonify({"error": "Invalid JSON body"}), 400
        project_id = body.get("project_id")
        question   = body.get("question", "").strip()
        if not project_id or not question:
            return jsonify({"error": "project_id and question required"}), 400

        projects = load_projects()
        if project_id not in projects:
            return jsonify({"error": f"Project '{project_id}' not found."}), 404

        project        = projects[project_id]
        context        = ""
        page_numbers   = []
        relevant_chunks = []

        try:
            vs = load_vectorstore(project_id)
            if vs is not None:
                docs = vs.similarity_search(question, k=5)
                context_parts = []
                for doc in docs:
                    page_num   = doc.metadata.get("page_number", "?")
                    chunk_text = doc.page_content
                    relevant_chunks.append({"text": chunk_text, "page": page_num, "company": doc.metadata.get("company_name", "")})
                    context_parts.append(f"[Page {page_num}]: {chunk_text}")
                    if page_num not in page_numbers:
                        page_numbers.append(page_num)
                context = "\n\n".join(context_parts)
        except Exception as vs_err:
            print(f"[WARN] Vectorstore failed, falling back: {vs_err}")

        if not context:
            data_file = project.get("data_file", "")
            if data_file and os.path.exists(data_file):
                with open(data_file, "r", encoding="utf-8") as f:
                    tree = json.load(f)
                context = "\n\n".join(
                    f"Section: {n.get('title','')}\n{n.get('brief','')}\n{n.get('summary','')}"
                    for n in tree.get("nodes", [])[:12]
                )

        if not context:
            return jsonify({"error": "No document content found."}), 404

        system = (
            f"You are an expert analyst for the RFP document: \"{project['doc_title']}\" "
            f"by {project.get('company_name', 'Unknown')}. "
            "Answer questions based ONLY on the provided document context. "
            "Be precise, helpful, and cite page numbers when available."
        )
        answer = groq_chat(
            f"Document context:\n{context}\n\nQuestion: {question}",
            model="llama-3.3-70b-versatile",
            system=system,
            max_tokens=1000,
            temperature=0.3,
        )

        return jsonify({
            "answer":          answer,
            "source_pages":    sorted(page_numbers),
            "relevant_chunks": relevant_chunks,
            "project_title":   project["doc_title"],
        })

    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": f"Server error: {str(e)}"}), 500


@app.route("/api/analytics-data", methods=["GET"])
def analytics_data():
    projects  = load_projects()
    proj_list = list(projects.values())

    sector_dist = {}
    for p in proj_list:
        cat = p.get("category", "General")
        sector_dist[cat] = sector_dist.get(cat, 0) + 1

    wins      = sum(1 for p in proj_list if p.get("status") == "win")
    losses    = sum(1 for p in proj_list if p.get("status") == "loss")
    active    = sum(1 for p in proj_list if p.get("status") == "active")
    completed = wins + losses
    total_emp = sum(p.get("employee_count", 0) for p in proj_list)

    from collections import defaultdict
    monthly = defaultdict(int)
    for p in proj_list:
        created = p.get("created_at", "")
        if created:
            try:
                monthly[created[:7]] += 1
            except Exception:
                pass

    return jsonify({
        "sector_distribution": sector_dist,
        "total_projects":      len(proj_list),
        "wins":                wins,
        "losses":              losses,
        "active_projects":     active,
        "completed_projects":  completed,
        "employees":           total_emp,
        "monthly_trend":       dict(sorted(monthly.items())[-6:]),
        "win_rate":            round((wins / completed * 100) if completed > 0 else 0, 1),
        "projects":            proj_list,
    })


@app.route("/api/rfp-intelligence", methods=["POST"])
def rfp_intelligence():
    try:
        if not os.getenv("GROQ_API_KEY"):
            return jsonify({"error": "GROQ_API_KEY is not set."}), 500

        full_text   = ""
        source_name = ""

        if request.is_json:
            body       = request.get_json()
            project_id = body.get("project_id")
            if not project_id:
                return jsonify({"error": "Provide 'project_id' in JSON body or upload a file."}), 400
            projects = load_projects()
            if project_id not in projects:
                return jsonify({"error": f"Project '{project_id}' not found."}), 404
            project = projects[project_id]
            try:
                pages     = get_pages_for_project(project)
                full_text = get_full_text_from_pages(pages)
            except FileNotFoundError:
                data_file = project.get("data_file", "")
                if data_file and os.path.exists(data_file):
                    with open(data_file, "r", encoding="utf-8") as fh:
                        tree = json.load(fh)
                    full_text = "\n\n".join(
                        f"{n.get('title','')}: {n.get('brief','')} {n.get('summary','')}"
                        for n in tree.get("nodes", [])
                    )
            source_name = project.get("doc_title", project_id)

        elif "file" in request.files:
            file = request.files["file"]
            ext  = os.path.splitext(file.filename)[1].lower()
            if ext not in (".pdf", ".txt"):
                return jsonify({"error": "Only PDF and TXT files are supported."}), 400
            tmp_dir  = "uploads"
            os.makedirs(tmp_dir, exist_ok=True)
            tmp_path = os.path.join(tmp_dir, f"intel_tmp_{uuid.uuid4().hex[:8]}{ext}")
            file.save(tmp_path)
            pages     = extract_text_with_pages(tmp_path, describe_images=True)
            full_text = get_full_text_from_pages(pages)
            source_name = file.filename
            try:
                os.remove(tmp_path)
            except Exception:
                pass
        else:
            return jsonify({"error": "Send a file or JSON body with 'project_id'."}), 400

        if not full_text.strip():
            return jsonify({"error": "Could not extract text from the document."}), 422

        llm          = GroqLLM()
        intelligence = extract_rfp_intelligence(llm, full_text)

        return jsonify({
            "source":          source_name,
            "extracted_at":    datetime.now().isoformat(),
            "domain":          intelligence.get("domain"),
            "sector":          intelligence.get("sector"),
            "summary":         intelligence.get("summary"),
            "keywords":        intelligence.get("keywords", []),
            "pricing":         intelligence.get("pricing", {}),
            "stakeholders":    intelligence.get("stakeholders", []),
            "deadlines":       intelligence.get("deadlines", []),
            "tech_categories": intelligence.get("tech_categories", []),
        })

    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": f"Server error: {str(e)}"}), 500


# ─── PPT GENERATION ────────────────────────────────────────────────────────────
PPT_SLIDES_JSON = "ppt_slides.json"
PPT_PPTX_NAME   = "presentation.pptx"

TITLE_SLIDE_LAYOUT = 0
MAIN_SLIDE_LAYOUT  = 1


def _ppt_get_ph_by_idx(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _ppt_build_semantic(full_text: str) -> dict:
    """Build slide semantic JSON via Groq (llama-3.3-70b-versatile)."""
    prompt = f"""You are an expert enterprise document analyst.

TASK: Analyze the document and produce a semantic JSON model for a PowerPoint presentation.

CONTENT RULES:
- Do NOT introduce new facts
- Rephrase content in professional, executive language
- Convert bullet-like text into cohesive narrative summaries
- Maintain factual fidelity

OUTPUT JSON FORMAT (strict — no markdown, no explanation):
{{
  "document_metadata": {{
    "title": "",
    "client": "",
    "document_type": ""
  }},
  "sections": [
    {{
      "heading": "",
      "section_number": "01",
      "summary": "",
      "key_points": ["point 1", "point 2", "point 3"],
      "signals": ["signal 1", "signal 2", "signal 3"],
      "sources": ["Source Name - Report Title Year"]
    }}
  ]
}}

RULES:
- Output MUST be valid JSON only
- Include ALL major sections (up to 12)
- section_number: zero-padded string like "01", "02" etc.
- summary: 2-3 sentence narrative per section
- key_points: 3-5 bullets per section, each under 20 words
- signals: 2-3 supporting signals or data points per section (short bullets)
- sources: 1-3 credible sources per section (real org names + report titles)
- Close all strings, arrays, and objects

DOCUMENT:
{full_text[:14000]}"""

    raw     = groq_chat(prompt, model="llama-3.3-70b-versatile", max_tokens=4000, temperature=0.2)
    cleaned = re.sub(r"```(?:json)?", "", raw).strip().rstrip("`").strip()
    m       = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if m:
        cleaned = m.group(0)
    return json.loads(cleaned)


def _ppt_create(data: dict, template_path: str, out_path: str) -> str:
    """Build PPTX from template preserving layout/styling."""
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.oxml.ns import qn

    def set_body(ph, summary, key_points):
        tf = ph.text_frame
        tf.clear()
        if summary:
            p0 = tf.paragraphs[0]
            p0.text = summary
            p0.font.size = Pt(11)
        for pt in (key_points or []):
            p = tf.add_paragraph()
            p.text = pt
            p.level = 1
            p.font.size = Pt(10)

    prs      = Presentation(template_path)
    meta     = data.get("document_metadata", {})
    sections = data.get("sections", [])

    title_layout   = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[MAIN_SLIDE_LAYOUT]
    try:
        thankyou_layout = prs.slide_layouts[16]
    except IndexError:
        thankyou_layout = prs.slide_layouts[-1]

    # Remove existing template slides
    xml_slides = prs.slides._sldIdLst
    slide_rids = [el.get(qn("r:id")) for el in list(xml_slides)]
    for rid in slide_rids:
        try:
            prs.part.drop_rel(rid)
        except Exception:
            pass
    for child in list(xml_slides):
        xml_slides.remove(child)

    # Title slide
    ts   = prs.slides.add_slide(title_layout)
    t_ph = _ppt_get_ph_by_idx(ts, 0)
    s_ph = _ppt_get_ph_by_idx(ts, 12)
    if t_ph: t_ph.text = meta.get("title", "RFP Presentation")
    if s_ph: s_ph.text = meta.get("client", "")

    # Content slides
    for section in sections:
        ns   = prs.slides.add_slide(content_layout)
        t_ph = _ppt_get_ph_by_idx(ns, 0)
        b_ph = _ppt_get_ph_by_idx(ns, 1) or _ppt_get_ph_by_idx(ns, 43)
        if t_ph: t_ph.text = section.get("heading", "")
        if b_ph: set_body(b_ph, section.get("summary", ""), section.get("key_points", []))

    # Thank-you slide
    ty    = prs.slides.add_slide(thankyou_layout)
    ty_ph = _ppt_get_ph_by_idx(ty, 15) or _ppt_get_ph_by_idx(ty, 0)
    if ty_ph:
        ty_ph.text = data.get("thankyou_name", "")

    prs.save(out_path)
    return out_path


@app.route("/api/ppt/generate", methods=["POST"])
def generate_ppt():
    try:
        if not os.getenv("GROQ_API_KEY"):
            return jsonify({"error": "GROQ_API_KEY is not set."}), 500

        body          = request.get_json(force=True) or {}
        project_id    = body.get("project_id", "").strip()
        template      = body.get("template", "CUSTOMPPT.pptx").strip()
        thankyou_name = body.get("thankyou_name", "")

        if not project_id:
            return jsonify({"error": "project_id is required"}), 400

        projects = load_projects()
        if project_id not in projects:
            return jsonify({"error": "Project not found"}), 404

        project   = projects[project_id]
        file_path = project.get("file_path", "")
        doc_title = project.get("doc_title", "RFP Presentation")

        if not file_path or not os.path.exists(file_path):
            return jsonify({"error": f"Source PDF not found: {file_path}"}), 404

        if not os.path.exists(template):
            return jsonify({"error": f"Template not found: '{template}'"}), 404

        try:
            pages = get_pages_for_project(project)
        except FileNotFoundError:
            pages = extract_text_with_pages(file_path, describe_images=True)
            save_structured_txt(pages, _project_folder(project))
        full_text = get_full_text_from_pages(pages)

        semantic                  = _ppt_build_semantic(full_text)
        semantic["thankyou_name"] = thankyou_name

        folder    = _project_folder(project)
        json_path = os.path.join(folder, PPT_SLIDES_JSON)
        ppt_json_payload = {
            "project_id":   project_id,
            "doc_title":    doc_title,
            "company_name": project.get("company_name", ""),
            "generated_at": datetime.now().isoformat(),
            "template":     template,
            **semantic,
        }
        with open(json_path, "w", encoding="utf-8") as jf:
            json.dump(ppt_json_payload, jf, indent=2, ensure_ascii=False)

        out_path = os.path.join(folder, PPT_PPTX_NAME)
        _ppt_create(semantic, template, out_path)

        sections    = semantic.get("sections", [])
        slide_count = len(sections) + 2
        actual_title = semantic.get("document_metadata", {}).get("title", doc_title)
        print(f"[INFO] PPT saved → {out_path}  ({slide_count} slides)")

        return jsonify({
            "success":      True,
            "download_url": f"/api/ppt/download/{project_id}",
            "json_url":     f"/api/ppt/json/{project_id}",
            "filename":     PPT_PPTX_NAME,
            "slide_count":  slide_count,
            "doc_title":    actual_title,
            "company_name": project.get("company_name", ""),
            "semantic":     semantic,
        })

    except json.JSONDecodeError as e:
        return jsonify({"error": f"AI returned invalid JSON — try again. Detail: {str(e)}"}), 500
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": f"Server error: {str(e)}"}), 500


@app.route("/api/ppt/json/<project_id>", methods=["GET"])
def get_ppt_json(project_id):
    projects = load_projects()
    if project_id not in projects:
        return jsonify({"error": "Project not found"}), 404
    folder = _project_folder(projects[project_id])
    path   = os.path.join(folder, PPT_SLIDES_JSON)
    if not os.path.exists(path):
        return jsonify({"error": "No PPT JSON found. Generate first."}), 404
    with open(path, "r") as f:
        return jsonify(json.load(f))


@app.route("/api/ppt/rebuild", methods=["POST"])
def rebuild_ppt():
    try:
        if not os.getenv("GROQ_API_KEY"):
            return jsonify({"error": "GROQ_API_KEY is not set."}), 500

        body       = request.get_json(force=True) or {}
        project_id = body.get("project_id", "").strip()
        template   = body.get("template", "CUSTOMPPT.pptx").strip()
        slides_in  = body.get("slides", [])

        if not project_id: return jsonify({"error": "project_id is required"}), 400
        if not slides_in:  return jsonify({"error": "slides array is required"}), 400
        if not os.path.exists(template): return jsonify({"error": f"Template not found: '{template}'"}), 404

        title_slide    = next((s for s in slides_in if s.get("type") == "title"), {})
        content_slides = [s for s in slides_in if s.get("type") == "content"]
        ty_slide       = next((s for s in slides_in if s.get("type") == "thankyou"), {})

        semantic = {
            "document_metadata": {
                "title":         title_slide.get("title",    "RFP Presentation"),
                "client":        title_slide.get("client",   ""),
                "document_type": title_slide.get("doc_type", ""),
            },
            "sections": [
                {
                    "heading":        s.get("heading", ""),
                    "section_number": s.get("section_number", str(i+1).zfill(2)),
                    "summary":        s.get("summary", ""),
                    "key_points":     s.get("key_points", []),
                    "signals":        s.get("signals", []),
                    "sources":        s.get("sources", []),
                }
                for i, s in enumerate(content_slides)
            ],
            "thankyou_name": ty_slide.get("name", ""),
        }

        all_projects = load_projects()
        folder       = _project_folder(all_projects.get(project_id, {}))
        out_path     = os.path.join(folder, PPT_PPTX_NAME)
        json_path    = os.path.join(folder, PPT_SLIDES_JSON)
        _ppt_create(semantic, template, out_path)

        if os.path.exists(json_path):
            with open(json_path, "r") as jf:
                existing = json.load(jf)
        else:
            existing = {"project_id": project_id}

        existing.update({
            "document_metadata": semantic["document_metadata"],
            "sections":          semantic["sections"],
            "thankyou_name":     semantic["thankyou_name"],
            "rebuilt_at":        datetime.now().isoformat(),
        })
        with open(json_path, "w") as jf:
            json.dump(existing, jf, indent=2)

        slide_count = len(content_slides) + 2
        return jsonify({"success": True, "download_url": f"/api/ppt/download/{project_id}", "slide_count": slide_count})

    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@app.route("/api/ppt/list", methods=["GET"])
def list_ppt_json():
    results  = []
    projects = load_projects()
    for pid, proj in projects.items():
        folder    = _project_folder(proj)
        json_path = os.path.join(folder, PPT_SLIDES_JSON)
        pptx_path = os.path.join(folder, PPT_PPTX_NAME)
        if not os.path.exists(json_path):
            continue
        try:
            with open(json_path) as f:
                data = json.load(f)
            results.append({
                "project_id":   pid,
                "folder":       proj.get("folder_name", pid),
                "doc_title":    data.get("doc_title", pid),
                "company_name": data.get("company_name", ""),
                "generated_at": data.get("generated_at", ""),
                "slide_count":  len(data.get("sections", [])) + 2,
                "has_pptx":     os.path.exists(pptx_path),
                "download_url": f"/api/ppt/download/{pid}",
                "json_url":     f"/api/ppt/json/{pid}",
            })
        except Exception:
            pass
    return jsonify(results)


@app.route("/api/ppt/download/<project_id>", methods=["GET"])
def download_ppt(project_id):
    from flask import send_file
    projects = load_projects()
    if project_id not in projects:
        return jsonify({"error": "Project not found"}), 404
    folder   = _project_folder(projects[project_id])
    ppt_path = os.path.join(folder, PPT_PPTX_NAME)
    if not os.path.exists(ppt_path):
        return jsonify({"error": "Not generated yet."}), 404
    folder_name   = projects[project_id].get("folder_name", project_id)
    download_name = f"{folder_name}.pptx"
    return send_file(
        os.path.abspath(ppt_path),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=download_name,
    )


# ─── MARKET RESEARCH ──────────────────────────────────────────────────────────
CARD_DEFS = [
    {"id":"market_analysis","title":"Market analysis","subtitle":"Category framing + demand signals","seed_terms":["market","industry","growth","demand","trend","strategy","objective","business","value","opportunity","scope","background"]},
    {"id":"competitive_analysis","title":"Competitive analysis","subtitle":"Competitive arena + evaluation criteria","seed_terms":["compet","vendor","supplier","evaluation","criteria","scoring","selection","qualification","capability","benchmark"]},
    {"id":"technology_adoption","title":"Technology adoption","subtitle":"Adoption enablers/risks + integration signals","seed_terms":["technology","technical","architecture","system","integration","data","platform","cloud","security","compliance","automation"]},
    {"id":"client_behaviour_study","title":"Client behaviour and study","subtitle":"Stakeholders + ways of working + change signals","seed_terms":["stakeholder","user","process","workflow","manual","pain point","requirement","governance","change","training","adoption","team"]},
    {"id":"regional_market_opportunities","title":"Regional market opportunities","subtitle":"Geo signals + rollout/onboarding implications","seed_terms":["region","country","global","local","market","geography","language","regulatory","compliance","expansion","rollout"]},
    {"id":"pricing_strategy_benchmarks","title":"Pricing strategy benchmarks","subtitle":"Commercial model signals + benchmarking template","seed_terms":["price","pricing","cost","budget","commercial","payment","licens","rate","fee","SLA","contract","terms","milestone"]},
]


def _mr_norm(text): return re.sub(r"\s+", " ", (text or "")).strip().lower()
def _mr_node_text_blob(node):
    parts = [node.get("id",""),node.get("title",""),node.get("desc",""),node.get("tag",""),node.get("summary",""),node.get("brief","")," ".join(node.get("keywords",[]) or [])," ".join(node.get("taglines",[]) or [])," ".join(node.get("suggestions",[]) or [])]
    return _mr_norm(" ".join([p for p in parts if p]))
def _mr_score_node(node_blob, seed_terms, overall_terms):
    score = 0.0
    for t in seed_terms:
        if _mr_norm(t) and _mr_norm(t) in node_blob: score += 3.0
    for t in overall_terms:
        if _mr_norm(t) and _mr_norm(t) in node_blob: score += 1.0
    if any(x in node_blob for x in ["evaluation","criteria","scoring"]): score += 0.5
    if any(x in node_blob for x in ["budget","pricing","cost","commercial"]): score += 0.5
    return score
def _mr_pick_nodes(nodes, seed_terms, overall_terms, top_k=6):
    blobs  = [(n, _mr_node_text_blob(n)) for n in nodes]
    scored = [(n, _mr_score_node(blob, seed_terms, overall_terms)) for n, blob in blobs]
    scored.sort(key=lambda x: x[1], reverse=True)
    picked = [n for n, s in scored if s > 0][:top_k]
    if picked: return picked
    fallback = [(n, sum(1 for t in overall_terms if _mr_norm(t) and _mr_norm(t) in blob)) for n, blob in blobs]
    fallback.sort(key=lambda x: x[1], reverse=True)
    return [n for n, _ in fallback][:max(3, top_k//2)]
def _mr_compact_nodes(nodes, max_brief=650):
    compact = []
    for n in nodes:
        brief = (n.get("brief") or "")
        if len(brief) > max_brief: brief = brief[:max_brief].rstrip() + "…"
        compact.append({"id":n.get("id"),"title":n.get("title"),"tag":n.get("tag"),"desc":n.get("desc",""),"summary":n.get("summary",""),"keywords":n.get("keywords",[]) or [],"brief":brief})
    return compact
def _mr_build_nodes_from_report(report):
    nodes = []
    def add(nid, title, brief, kws=None, tag=""):
        nodes.append({"id":nid,"title":title,"tag":tag,"desc":"","summary":"","brief":brief,"keywords":kws or [],"taglines":[],"suggestions":[]})
    if report.get("summary"): add("rfp_summary","RFP summary",str(report["summary"]),keywords=report.get("keywords",[]),tag="summary")
    domain,sector = report.get("domain"),report.get("sector")
    if domain or sector: add("rfp_domain_sector","Domain & sector",f"Domain: {domain or ''} | Sector: {sector or ''}",keywords=[x for x in [domain,sector] if x],tag="context")
    kws = report.get("keywords",[])
    if kws: add("rfp_keywords","RFP keywords","Keywords: "+", ".join(kws[:40]),keywords=kws,tag="keywords")
    pricing = report.get("pricing") or {}
    if pricing: add("rfp_pricing","Commercial / pricing signals",f"Budget mentioned: {pricing.get('budget_mentioned')} | Estimated value: {pricing.get('estimated_value')} | Model: {pricing.get('pricing_model')} | Currency: {pricing.get('currency')} | Notes: {pricing.get('notes')}",keywords=["pricing","commercial","budget"],tag="commercial")
    deadlines = report.get("deadlines",[])
    if deadlines:
        lines = [f"{d.get('event')}: {d.get('date')}" for d in deadlines]
        add("rfp_deadlines","Timeline / deadlines"," | ".join(lines),keywords=["deadline","timeline"],tag="timeline")
    stakeholders = report.get("stakeholders",[])
    if stakeholders:
        lines = [f"{s.get('name')} ({s.get('role')})" for s in stakeholders]
        add("rfp_stakeholders","Stakeholders","Stakeholders: "+" | ".join(lines),keywords=["stakeholder","client"],tag="stakeholders")
    tech = report.get("tech_categories",[])
    if tech:
        lines = [f"{t.get('category')}: {', '.join(t.get('specifics',[]))}" for t in tech]
        add("rfp_tech","Technology scope"," | ".join(lines),keywords=["technology","data","cloud"],tag="technology")
    if not nodes: add("rfp_fallback","RFP content","No structured fields found.",tag="fallback")
    return nodes


def _mr_generate_card_phase1(groq_client: Groq, model: str, card_def: dict, meta: dict, evidence_nodes: list) -> dict:
    system_msg = (
        "You are a market research analyst generating Phase-1 RFP-grounded insights.\n"
        "Use ONLY the evidence provided from the RFP extraction.\n"
        "Output must be valid JSON matching the schema exactly.\n"
        "Schema: {id, title, subtitle, tags: [], snapshot: [], rfp_insights: [{insight, evidence_node_ids}], assumptions: [], research_questions: [], external_facts: []}"
    )
    user_payload = {
        "meta": {"doc_title":meta.get("doc_title",""),"company_name":meta.get("company_name",""),"sector":meta.get("sector",""),"domain":meta.get("domain",""),"primary_keywords":meta.get("primary_keywords",[])},
        "card": {"id":card_def["id"],"title":card_def["title"],"subtitle":card_def["subtitle"]},
        "instructions": {"snapshot_rules":"Write 3-5 crisp bullets, each < 20 words.","rfp_insights_rules":"Write 3-5 insights supported by evidence_node_ids.","assumptions_rules":"Write 2-3 plausible assumptions.","research_questions_rules":"Write 2-4 sharp validation questions.","tags_rules":"Write 3-5 short tags.","external_facts_rules":"Leave external_facts as empty list."},
        "evidence_nodes": _mr_compact_nodes(evidence_nodes),
    }
    response  = groq_client.chat.completions.create(model=model, messages=[{"role":"system","content":system_msg},{"role":"user","content":json.dumps(user_payload)}], temperature=0.3, max_tokens=2000)
    raw       = response.choices[0].message.content
    card_data = parse_json_response(raw)
    card_data["id"]       = card_def["id"]
    card_data["title"]    = card_def["title"]
    card_data["subtitle"] = card_def["subtitle"]
    card_data.setdefault("tags", [])
    card_data.setdefault("snapshot", [])
    card_data.setdefault("rfp_insights", [])
    card_data.setdefault("assumptions", [])
    card_data.setdefault("research_questions", [])
    card_data["external_facts"] = []
    return card_data


_SOURCE_HOMEPAGES: Dict[str, str] = {
    "gartner":"https://www.gartner.com/en","mckinsey":"https://www.mckinsey.com","mckinsey & company":"https://www.mckinsey.com","statista":"https://www.statista.com","idc":"https://www.idc.com","forrester":"https://www.forrester.com","deloitte":"https://www.deloitte.com","pwc":"https://www.pwc.com","accenture":"https://www.accenture.com","bcg":"https://www.bcg.com","boston consulting group":"https://www.bcg.com","world bank":"https://www.worldbank.org","ibm":"https://www.ibm.com","harvard business review":"https://hbr.org","hbr":"https://hbr.org","bain":"https://www.bain.com","kpmg":"https://home.kpmg","ey":"https://www.ey.com","oecd":"https://www.oecd.org","imf":"https://www.imf.org","wef":"https://www.weforum.org","world economic forum":"https://www.weforum.org","bloomberg":"https://www.bloomberg.com","reuters":"https://www.reuters.com",
}

def _get_source_homepage(source_name):
    key = (source_name or "").lower().strip()
    if key in _SOURCE_HOMEPAGES: return _SOURCE_HOMEPAGES[key]
    for k, v in _SOURCE_HOMEPAGES.items():
        if k in key or key in k: return v
    return ""


def _mr_enrich_card_phase2(groq_client: Groq, model: str, meta: dict, card: dict, max_facts: int = 5) -> dict:
    system_msg = (
        "You are a senior market research analyst. Generate real, verifiable external facts.\n"
        "For each fact provide: claim (specific, data-rich), source_name (real publisher), url (real top-level URL), confidence (high/medium/low).\n"
        "Output ONLY valid JSON: { \"external_facts\": [{claim, source_name, url, confidence}] }"
    )
    user_payload = {
        "context": {"sector":meta.get("sector",""),"domain":meta.get("domain",""),"doc_title":meta.get("doc_title","")},
        "card": {"title":card.get("title"),"research_questions":card.get("research_questions",[]),"assumptions":card.get("assumptions",[])},
        "instructions": f"Generate {max_facts} distinct external facts for '{card.get('title')}'. Prefer: Gartner, McKinsey, Statista, IDC, Forrester, Deloitte, BCG, WEF.",
    }
    try:
        response = groq_client.chat.completions.create(model=model, messages=[{"role":"system","content":system_msg},{"role":"user","content":json.dumps(user_payload)}], temperature=0.25, max_tokens=1800)
        raw      = response.choices[0].message.content
        data     = parse_json_response(raw)
        facts    = (data.get("external_facts") or [])[:max_facts]
        for fact in facts:
            homepage         = _get_source_homepage(fact.get("source_name",""))
            fact["homepage"] = homepage
            if not fact.get("url"): fact["url"] = homepage or ""
        card["external_facts"] = facts
    except Exception as e:
        print(f"[WARN] Phase 2 failed for {card.get('id')}: {e}")
        card["external_facts"] = []
    return card


def generate_market_research_for_project(project_id: str, report_data: dict, data_file: str = None, enable_web: bool = True, model: str = "llama-3.3-70b-versatile") -> str:
    groq_client = Groq(api_key=get_api_key())
    nodes = []
    if data_file and os.path.exists(data_file):
        try:
            with open(data_file) as f:
                data_json = json.load(f)
            nodes = data_json.get("nodes", [])
        except Exception as e:
            print(f"[WARN] Could not load data_file {data_file}: {e}")
    if not nodes:
        nodes = _mr_build_nodes_from_report(report_data)

    meta = {
        "doc_title":          report_data.get("doc_title",""),
        "company_name":       report_data.get("company_name",""),
        "sector":             report_data.get("sector",""),
        "domain":             report_data.get("domain",""),
        "primary_keywords":   report_data.get("keywords",[]) or [],
        "supporting_keywords": [],
    }
    overall_terms = meta["primary_keywords"][:40]

    cards = []
    for cd in CARD_DEFS:
        ev   = _mr_pick_nodes(nodes, cd["seed_terms"], overall_terms, top_k=6)
        print(f"[MR] Phase 1: {cd['id']} (evidence={len(ev)})", flush=True)
        card = _mr_generate_card_phase1(groq_client, model, cd, meta, ev)
        if enable_web:
            print(f"[MR] Phase 2: {cd['id']}", flush=True)
            card = _mr_enrich_card_phase2(groq_client, model, meta, card, max_facts=5)
        cards.append(card)

    doc      = {"meta": meta, "cards": cards}
    all_projects = load_projects()
    proj     = all_projects.get(project_id, {})
    folder   = _project_folder(proj)
    out_path = os.path.join(folder, "market_research.json")
    with open(out_path, "w") as f:
        json.dump(doc, f, indent=2, ensure_ascii=False)
    print(f"[MR] Saved → {out_path}")
    return out_path


@app.route("/api/market-research/projects", methods=["GET"])
def mr_list_projects():
    projects = load_projects()
    result   = []
    for pid, proj in projects.items():
        folder     = _project_folder(proj)
        intel_path = os.path.join(folder, "intelligence.json")
        mr_file    = os.path.join(folder, "market_research.json")
        result.append({
            "id":                  pid,
            "doc_title":           proj.get("doc_title",""),
            "company_name":        proj.get("company_name",""),
            "sector":              proj.get("sector") or proj.get("category",""),
            "domain":              proj.get("domain",""),
            "has_report":          os.path.exists(intel_path),
            "has_market_research": os.path.exists(mr_file),
            "created_at":          proj.get("created_at",""),
        })
    return jsonify(result)


@app.route("/api/market-research/generate", methods=["POST"])
def mr_generate():
    try:
        body       = request.get_json()
        project_id = body.get("project_id")
        enable_web = body.get("enable_web", True)
        if not project_id: return jsonify({"error":"project_id required"}), 400
        projects = load_projects()
        if project_id not in projects: return jsonify({"error":"Project not found"}), 404
        proj        = projects[project_id]
        folder      = _project_folder(proj)
        intel_path  = os.path.join(folder, "intelligence.json")
        report_data = None
        if os.path.exists(intel_path):
            with open(intel_path) as f:
                report_data = json.load(f)
        if not report_data:
            report_data = {"doc_title":proj.get("doc_title",""),"company_name":proj.get("company_name",""),"sector":proj.get("sector") or proj.get("category",""),"domain":proj.get("domain",""),"keywords":[]}
        out_path = generate_market_research_for_project(project_id, report_data, data_file=proj.get("data_file"), enable_web=enable_web)
        return jsonify({"success":True,"project_id":project_id,"path":out_path})
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error":str(e)}), 500


@app.route("/api/market-research/<project_id>", methods=["GET"])
def mr_get(project_id):
    projects = load_projects()
    if project_id not in projects: return jsonify({"error":"Project not found"}), 404
    folder  = _project_folder(projects[project_id])
    mr_file = os.path.join(folder, "market_research.json")
    if not os.path.exists(mr_file): return jsonify({"error":"Market research not yet generated."}), 404
    with open(mr_file) as f: data = json.load(f)
    return jsonify(data)


@app.route("/api/market-research/report/<project_id>", methods=["GET"])
def mr_get_report(project_id):
    projects = load_projects()
    if project_id not in projects: return jsonify({"error":"Project not found"}), 404
    folder     = _project_folder(projects[project_id])
    intel_path = os.path.join(folder, "intelligence.json")
    if not os.path.exists(intel_path): return jsonify({"error":"No intelligence report found."}), 404
    with open(intel_path) as f: data = json.load(f)
    return jsonify(data)


if __name__ == "__main__":
    app.run(debug=True, port=5000)